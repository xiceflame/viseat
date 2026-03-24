import os

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import font_manager
from pptx import Presentation
from pptx.util import Inches, Pt


MODELS = ["智谱清言", "Gemini", "Poe", "通义千问", "豆包", "ChatGPT", "Claude", "Perplexity", "Mistral Le Chat", "Kimi"]
F1 = np.array([56.0, 50.0, 48.3, 45.0, 44.7, 43.0, 37.1, 34.9, 25.8, 19.2])
WEIGHT_ERR = np.array([54.9, 25.8, 37.0, 27.7, 44.9, 58.0, 40.6, 34.7, 33.1, 51.5])
NAME_ACC = np.array([60, 80, 60, 70, 80, 50, 70, 80, 60, 50])

TONGYI_NAME = "通义千问"
TONGYI_COLOR = "#FF8A00"
NEUTRAL_COLOR = "#A7A9AD"

OUTDIR = "viseat_output"


def setup_fonts() -> None:
    candidates = ["Hiragino Sans GB", "STHeiti", "Songti SC"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    chosen = next((n for n in candidates if n in available), None)
    if chosen:
        mpl.rcParams["font.family"] = "sans-serif"
        mpl.rcParams["font.sans-serif"] = [chosen]
    mpl.rcParams["axes.unicode_minus"] = False


def save_fig(path: str) -> None:
    plt.savefig(path, bbox_inches="tight", dpi=200)
    plt.close()


def ensure_outdir() -> None:
    os.makedirs(OUTDIR, exist_ok=True)


def make_comparison_table_png() -> None:
    df = pd.DataFrame(
        {
            "模型": MODELS,
            "名称准确率(%)": NAME_ACC.astype(int),
            "重量平均误差(%)": WEIGHT_ERR,
            "食材拆分 F1(%)": F1,
        }
    )

    fig, ax = plt.subplots(figsize=(9, 1.6))
    ax.axis("off")

    table = ax.table(
        cellText=df.values.tolist(),
        colLabels=list(df.columns),
        colColours=["#F0F4FA"] * len(df.columns),
        cellLoc="center",
        loc="center",
    )

    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 1.2)

    tongyi_row_idx = MODELS.index(TONGYI_NAME) + 1
    for col_idx in range(len(df.columns)):
        table[tongyi_row_idx, col_idx].set_facecolor("#FFF3E0")

    save_fig(os.path.join(OUTDIR, "comparison_table.png"))


def make_radar_png() -> None:
    labels = ["名称(%)", "重量(逆向得分)", "食材F1(%)"]
    axes_count = len(labels)
    weight_score = 100 - WEIGHT_ERR

    angles = np.linspace(0, 2 * np.pi, axes_count, endpoint=False).tolist()
    angles += angles[:1]

    fig = plt.figure(figsize=(8, 6))
    ax = fig.add_subplot(111, polar=True)

    for i, model in enumerate(MODELS):
        row = [NAME_ACC[i], weight_score[i], F1[i]]
        values = row + row[:1]
        if model == TONGYI_NAME:
            ax.plot(angles, values, linewidth=2.0, label=model, color=TONGYI_COLOR)
            ax.fill(angles, values, alpha=0.12, color=TONGYI_COLOR)
        else:
            ax.plot(angles, values, linewidth=0.8, label=None, color=NEUTRAL_COLOR, alpha=0.9)

    ax.set_thetagrids(np.degrees(angles[:-1]), labels)
    ax.set_ylim(0, 100)
    ax.set_title("三维性能雷达（名称 / 重量(反向) / 食材F1）", y=1.08)
    ax.legend(loc="upper right")

    save_fig(os.path.join(OUTDIR, "radar_all_models.png"))


def make_composite_bar_png(weight_w: float = 0.6, f1_w: float = 0.25, name_w: float = 0.15) -> None:
    if abs(weight_w + f1_w + name_w - 1.0) > 1e-9:
        raise ValueError("weights must sum to 1.0")

    weight_score = 100 - WEIGHT_ERR
    name_norm = NAME_ACC / 100.0
    f1_norm = F1 / 100.0
    w_inv_norm = (weight_score - weight_score.min()) / (weight_score.max() - weight_score.min())

    composite = weight_w * w_inv_norm + f1_w * f1_norm + name_w * name_norm
    order = np.argsort(-composite)

    fig, ax = plt.subplots(figsize=(9, 3.0))
    bars = ax.bar(np.array(MODELS)[order], composite[order], color="#6A7BD8")

    for idx, rect in enumerate(bars):
        model_name = np.array(MODELS)[order][idx]
        if model_name == TONGYI_NAME:
            rect.set_color(TONGYI_COLOR)
            rect.set_edgecolor("black")
            rect.set_linewidth(1.2)

    ax.set_ylabel("综合得分（已归一化）")
    ax.set_title("加权复合得分（重量60% / F125% / 名称15%）")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()

    save_fig(os.path.join(OUTDIR, "composite_bar.png"))


def make_risk_quadrant_png() -> None:
    weight_score = 100 - WEIGHT_ERR
    weight_norm = (weight_score - weight_score.min()) / (weight_score.max() - weight_score.min())
    f1_norm = F1 / 100.0

    fig, ax = plt.subplots(figsize=(6, 5))
    ax.scatter(weight_norm, f1_norm, s=60, c=NEUTRAL_COLOR)

    tongyi_idx = MODELS.index(TONGYI_NAME)
    ax.scatter(
        weight_norm[tongyi_idx],
        f1_norm[tongyi_idx],
        s=180,
        c=TONGYI_COLOR,
        edgecolors="black",
    )
    ax.annotate(TONGYI_NAME, (weight_norm[tongyi_idx] + 0.02, f1_norm[tongyi_idx]))

    ax.set_xlabel("重量稳定性（越右越稳）")
    ax.set_ylabel("食材拆分质量 F1（越高越好）")
    ax.set_title("风险象限图：重量稳定性 vs 食材拆分质量")
    ax.axvline(0.5, color="#DADDE0")
    ax.axhline(0.5, color="#DADDE0")

    save_fig(os.path.join(OUTDIR, "risk_quadrant.png"))


def make_pptx() -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.5), Inches(0.6))
    tf = title.text_frame
    tf.text = "为什么选择「通义千问」作为 VisEat 首选模型"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Hiragino Sans GB"

    subtitle = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(9.5), Inches(0.3))
    subtitle.text_frame.text = "场景：营养助手（重量估算优先） · 评测日：2026-01-13"
    subtitle.text_frame.paragraphs[0].font.size = Pt(12)
    subtitle.text_frame.paragraphs[0].font.name = "Hiragino Sans GB"

    slide.shapes.add_picture(os.path.join(OUTDIR, "comparison_table.png"), Inches(0.4), Inches(1.5), width=Inches(6.8))
    slide.shapes.add_picture(os.path.join(OUTDIR, "composite_bar.png"), Inches(0.4), Inches(3.1), width=Inches(6.8))
    slide.shapes.add_picture(os.path.join(OUTDIR, "radar_all_models.png"), Inches(7.0), Inches(1.5), width=Inches(3.5))
    slide.shapes.add_picture(os.path.join(OUTDIR, "risk_quadrant.png"), Inches(7.0), Inches(3.8), width=Inches(3.5))

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(5.0), Inches(10.2), Inches(1.8))
    t = tb.text_frame
    t.text = (
        "方法学（摘要）：样本 10 道菜；指标 = 名称(严格)、重量误差(%)、食材 F1。\n"
        "纠偏建议（高优先级）：1) 名称规范化  2) 重量校准  3) 调料降权\n"
        "上线策略：默认通义千问，低置信度回退 Gemini / 人工复核；先行 A/B 小流量验证。"
    )
    for p in t.paragraphs:
        p.font.size = Pt(11)
        p.font.name = "Hiragino Sans GB"

    out_pptx = os.path.join(OUTDIR, "VisEat_Tongyi_choice_final.pptx")
    prs.save(out_pptx)
    return out_pptx


def main() -> None:
    setup_fonts()
    ensure_outdir()
    make_comparison_table_png()
    make_radar_png()
    make_composite_bar_png()
    make_risk_quadrant_png()
    out_pptx = make_pptx()

    print("Outputs written to:", OUTDIR)
    print("-", os.path.join(OUTDIR, "comparison_table.png"))
    print("-", os.path.join(OUTDIR, "radar_all_models.png"))
    print("-", os.path.join(OUTDIR, "composite_bar.png"))
    print("-", os.path.join(OUTDIR, "risk_quadrant.png"))
    print("-", out_pptx)


if __name__ == "__main__":
    main()
